#!/usr/bin/env python3
from __future__ import annotations

import argparse
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def render_pdf_pages(pdf_path: Path, out_dir: Path, dpi: int) -> list[Path]:
    prefix = out_dir / "slide"
    subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf_path), str(prefix)],
        check=True,
    )
    images = sorted(out_dir.glob("slide-*.png"))
    if not images:
        raise RuntimeError("No PNG pages were generated from the PDF.")
    return images


def build_pptx(images: list[Path], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Remove the default first slide so every slide comes from the PDF pages.
    while prs.slides:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    for image_path in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convert a PDF deck into a PPTX deck with one full-slide image per page."
    )
    parser.add_argument("pdf", type=Path, help="Source PDF path")
    parser.add_argument("pptx", type=Path, help="Target PPTX path")
    parser.add_argument("--dpi", type=int, default=180, help="Rasterization DPI")
    args = parser.parse_args()

    with tempfile.TemporaryDirectory(prefix="pdf_to_pptx_") as tmp:
        tmp_dir = Path(tmp)
        images = render_pdf_pages(args.pdf, tmp_dir, args.dpi)
        build_pptx(images, args.pptx)


if __name__ == "__main__":
    main()
